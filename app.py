import os
import faiss
import numpy as np
import pandas as pd
from langchain.document_loaders import PyPDFLoader
from sentence_transformers import SentenceTransformer
import streamlit as st
import tempfile
import openai
from dotenv import load_dotenv 
from docx import Document  
import json  
from pptx import Presentation 
import io  # Import io for handling file download

# Load environment variables
load_dotenv()

# Set up OpenAI API credentials
openai.api_type = "azure"
openai.api_version = os.getenv("AZURE_OPENAI_API_VERSION", "").strip()
API_KEY = os.getenv("AZURE_OPENAI_API_KEY", "").strip()
openai.api_key = API_KEY
RESOURCE_ENDPOINT = os.getenv("OPENAI_API_BASE", "").strip()
openai.api_base = RESOURCE_ENDPOINT
deployment = "TG-OAi-GPTModel" 

# Load Documents
def load_pdf(uploaded_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name
    loader = PyPDFLoader(temp_file_path)
    documents = loader.load_and_split()
    return [doc.page_content for doc in documents]

def load_csv(uploaded_file):
    df = pd.read_csv(uploaded_file)
    return df.astype(str).agg(' '.join, axis=1).tolist()  # Concatenate all columns into a single string per row

def load_json(uploaded_file):
    data = json.load(uploaded_file)
    if isinstance(data, list):
        return [json.dumps(item) for item in data]  # Return each item as a JSON string if it's a list
    else:
        return [json.dumps(data)]  # Return a single JSON object as a string

def load_docx(uploaded_file):
    doc = Document(uploaded_file)
    return [para.text for para in doc.paragraphs if para.text]

def load_xlsx(uploaded_file):
    df = pd.read_excel(uploaded_file)
    return df.astype(str).agg(' '.join, axis=1).tolist()  # Concatenate all columns into a single string per row

def load_pptx(uploaded_file):
    prs = Presentation(uploaded_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return text

# Embedding
model = SentenceTransformer('all-MiniLM-L6-v2')

def generate_embeddings(documents):
    texts = [doc for doc in documents]
    embeddings = model.encode(texts)
    return embeddings, texts

# Vector Stores
def create_faiss_index(embeddings):
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)
    return index

# Search function
def search_faiss(query, index, documents):
    query_embedding = model.encode([query])[0]
    distances, indices = index.search(np.array([query_embedding]), k=5)  # Search top 5
    results = [documents[i] for i in indices[0]]
    return results

# Get response from Azure OpenAI
def get_response(question, context):
    text_prompt = f"Given the following context:\n\n{context}\n\nPlease answer the question clearly and accurately in detail:\n\nQuestion: {question}\n"
    response = openai.ChatCompletion.create(
        engine=deployment,
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": text_prompt},
        ]
    )
    return response['choices'][0]['message']['content']

# Streamlit Interface
st.title("Document RAG Chatbot")

# Initialize session state for history if it doesn't exist
if 'history' not in st.session_state:
    st.session_state.history = []

# File uploader
uploaded_files = st.file_uploader("Upload your Documents here: ", accept_multiple_files=True, type=["pdf", "csv", "json", "docx", "xlsx", "pptx"])

if uploaded_files:
    documents = []
    for file in uploaded_files:
        if file.name.endswith('.pdf'):
            documents.extend(load_pdf(file))
        elif file.name.endswith('.csv'):
            documents.extend(load_csv(file))
        elif file.name.endswith('.json'):
            documents.extend(load_json(file))
        elif file.name.endswith('.docx'):
            documents.extend(load_docx(file))
        elif file.name.endswith('.xlsx'):
            documents.extend(load_xlsx(file))
        elif file.name.endswith('.pptx'):
            documents.extend(load_pptx(file))

    # Generate embeddings and create FAISS index
    if documents:  # Check if documents are loaded
        embeddings, texts = generate_embeddings(documents)
        faiss_index = create_faiss_index(np.array(embeddings))
    
        # User query
        query = st.text_input("Ask a question:")
    
        if query:
            # Search and generate answer
            retrieved_docs = search_faiss(query, faiss_index, documents)
            context = " ".join(retrieved_docs)
            answer = get_response(query, context)
            
            # Append to history
            st.session_state.history.append({"question": query, "answer": answer})
            st.write(f"Answer: {answer}")
        else:
            st.write("Please ask a question.")

# Button to see history
if st.button("See History"):
    if st.session_state.history:
        for entry in st.session_state.history:
            st.write(f"**Question:** {entry['question']}")
            st.write(f"**Answer:** {entry['answer']}")
    else:
        st.write("No history available.")

# Button to download history
def download_history():
    history_str = "\n".join([f"Question: {entry['question']}\nAnswer: {entry['answer']}" for entry in st.session_state.history])
    return history_str

if st.button("Download History"):
    history_text = download_history()
    st.download_button("Download", history_text, file_name="chat_history.txt")
